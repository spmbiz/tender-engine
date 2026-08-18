from __future__ import annotations

import argparse
import hashlib
import json
import os
import shutil
import subprocess
import tarfile
import zipfile
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def unpack_archive(path: Path) -> Path | None:
    out = path.parent / f"{path.stem}_unpacked"
    if out.exists() and any(out.iterdir()):
        return out
    out.mkdir(parents=True, exist_ok=True)
    try:
        if zipfile.is_zipfile(path):
            with zipfile.ZipFile(path) as z:
                z.extractall(out)
            return out
        if tarfile.is_tarfile(path):
            with tarfile.open(path) as t:
                t.extractall(out)
            return out
        if path.suffix.lower() == ".7z" and shutil.which("7z"):
            subprocess.run(["7z", "x", "-y", str(path), f"-o{out}"], check=False, capture_output=True)
            if any(out.iterdir()):
                return out
    except Exception:
        return None
    try:
        out.rmdir()
    except Exception:
        pass
    return None


def recursive_unpack(files_root: Path, max_rounds: int = 5):
    seen = set()
    for _ in range(max_rounds):
        changed = False
        for p in list(files_root.rglob("*")):
            if not p.is_file() or p in seen:
                continue
            seen.add(p)
            if p.suffix.lower() in {".zip", ".7z", ".tar", ".gz", ".tgz"} or zipfile.is_zipfile(p):
                out = unpack_archive(p)
                if out:
                    changed = True
        if not changed:
            break


def extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    except Exception as exc:
        return f"[PDF_EXTRACTION_ERROR {exc!r}]"


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        chunks = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                chunks.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(chunks)
    except Exception as exc:
        return f"[DOCX_EXTRACTION_ERROR {exc!r}]"


def extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        chunks = []
        for ws in wb.worksheets:
            chunks.append(f"[SHEET {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v) for v in row]
                if any(vals):
                    chunks.append(" | ".join(vals))
        return "\n".join(chunks)
    except Exception as exc:
        return f"[XLSX_EXTRACTION_ERROR {exc!r}]"


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        chunks = []
        for idx, slide in enumerate(prs.slides, 1):
            chunks.append(f"[SLIDE {idx}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks)
    except Exception as exc:
        return f"[PPTX_EXTRACTION_ERROR {exc!r}]"


def extract_text(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return "pdf", extract_pdf(path)
    if ext == ".docx":
        return "docx", extract_docx(path)
    if ext in {".xlsx", ".xlsm"}:
        return "xlsx", extract_xlsx(path)
    if ext == ".pptx":
        return "pptx", extract_pptx(path)
    if ext in {".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".rtf"}:
        try:
            return "text", path.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            return "text", f"[TEXT_EXTRACTION_ERROR {exc!r}]"
    return "unsupported", ""


def _native_text_yield(text: str) -> int:
    if not text or "_EXTRACTION_ERROR " in text[:80]:
        return 0
    return len(text.strip())


def _load_docling_adapter():
    try:
        from pipeline.docling_deep_parser import parse_with_docling, should_deep_parse
        return parse_with_docling, should_deep_parse
    except ModuleNotFoundError as first:
        if getattr(first, "name", None) != "pipeline":
            raise
        from docling_deep_parser import parse_with_docling, should_deep_parse
        return parse_with_docling, should_deep_parse


def _try_docling(path: Path, root: Path, kind: str, text: str) -> tuple[str, dict | None]:
    mode = os.getenv("DCE_DOCLING_MODE", "auto").strip().lower()
    min_chars = max(0, int(os.getenv("DCE_DOCLING_MIN_CHARS", "800")))
    try:
        parse_with_docling, should_deep_parse = _load_docling_adapter()
    except Exception as exc:
        return text, {"status": "ADAPTER_IMPORT_ERROR", "error": repr(exc)}

    if not should_deep_parse(kind, _native_text_yield(text), mode=mode, min_chars=min_chars):
        return text, None

    deep = parse_with_docling(path, root / "derived" / "docling")
    if deep.get("status") != "DOWNLOADED_SOURCE_PARSED":
        return text, deep

    try:
        deep_text = Path(str(deep["markdown_path"])).read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        deep["status"] = "SIDECAR_READ_ERROR"
        deep["error"] = repr(exc)
        return text, deep

    if len(deep_text.strip()) > _native_text_yield(text):
        return deep_text, deep
    return text, deep


def _load_manifest(root: Path) -> dict:
    try:
        obj = json.loads((root / "manifest.json").read_text(encoding="utf-8", errors="replace"))
        return obj if isinstance(obj, dict) else {}
    except Exception:
        return {}


def _manifest_file_indexes(manifest: dict):
    """Map persisted file provenance without guessing across unrelated files."""
    by_sha: dict[str, dict] = {}
    by_name: dict[str, list[dict]] = {}
    for row in manifest.get("files") or []:
        if not isinstance(row, dict):
            continue
        sha = str(row.get("sha256") or row.get("sha") or "").strip().lower()
        if sha:
            by_sha[sha] = row
        for value in (row.get("name"), row.get("filename"), row.get("path")):
            name = Path(str(value or "")).name.casefold()
            if name:
                by_name.setdefault(name, []).append(row)
    return by_sha, by_name


def _source_provenance(path: Path, digest: str, by_sha: dict[str, dict], by_name: dict[str, list[dict]]) -> dict:
    row = by_sha.get(digest.lower())
    match_method = "SHA256" if row else None
    if row is None:
        candidates = by_name.get(path.name.casefold()) or []
        if len(candidates) == 1:
            row = candidates[0]
            match_method = "UNIQUE_FILENAME"
    row = row or {}
    return {
        "source_url": row.get("source_url") or row.get("url"),
        "source_manifest_name": row.get("name") or row.get("filename") or row.get("path"),
        "source_provenance_match": match_method,
    }


def process_candidate(root: Path):
    files_root = root / "files"
    if not files_root.exists():
        return None
    recursive_unpack(files_root)
    manifest = _load_manifest(root)
    by_sha, by_name = _manifest_file_indexes(manifest)
    docs = []
    corpus_chunks: list[str] = []
    corpus_cursor = 0
    for path in sorted(files_root.rglob("*")):
        if not path.is_file():
            continue
        kind, native_text = extract_text(path)
        text, deep = _try_docling(path, root, kind, native_text)
        rel = str(path.relative_to(root))
        digest = sha256_file(path)
        if text and len(text) > 1_500_000:
            text = text[:1_500_000] + "\n[TRUNCATED_AT_1_500_000_CHARS]"

        header = f"\n===== FILE: {rel} =====\n" if text else ""
        if header:
            corpus_chunks.append(header)
            corpus_cursor += len(header)
        text_start = corpus_cursor if text else None
        if text:
            corpus_chunks.append(text)
            corpus_cursor += len(text)
        text_end = corpus_cursor if text else None
        if text:
            corpus_chunks.append("\n")
            corpus_cursor += 1

        rec = {
            "path": rel,
            "name": path.name,
            "size": path.stat().st_size,
            "sha256": digest,
            "kind": kind,
            "native_text_chars": _native_text_yield(native_text),
            "text_chars": len(text),
            "corpus_start": text_start,
            "corpus_end": text_end,
            **_source_provenance(path, digest, by_sha, by_name),
        }
        if deep is not None:
            rec["deep_parser"] = deep
        docs.append(rec)

    corpus = "".join(corpus_chunks)
    (root / "document_index.json").write_text(json.dumps(docs, indent=2, ensure_ascii=False), encoding="utf-8")
    (root / "corpus.txt").write_text(corpus, encoding="utf-8")
    return {
        "root": str(root),
        "documents": len(docs),
        "documents_with_source_url": sum(1 for d in docs if d.get("source_url")),
        "documents_with_corpus_offsets": sum(1 for d in docs if d.get("corpus_start") is not None),
        "corpus_chars": len(corpus),
    }


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--root", default="out")
    ap.add_argument("--workers", type=int, default=int(os.getenv("DCE_EXTRACT_WORKERS", "2")))
    args = ap.parse_args()
    base = Path(args.root)
    candidates = sorted(set(p.parent for p in base.rglob("manifest.json")))
    workers = max(1, min(args.workers, len(candidates) or 1))
    results = []

    if workers == 1:
        for root in candidates:
            rec = process_candidate(root)
            if rec:
                results.append(rec)
    else:
        with ProcessPoolExecutor(max_workers=workers) as pool:
            futs = {pool.submit(process_candidate, root): root for root in candidates}
            for fut in as_completed(futs):
                try:
                    rec = fut.result()
                except Exception as exc:
                    rec = {"root": str(futs[fut]), "error": repr(exc), "documents": 0, "corpus_chars": 0}
                if rec:
                    results.append(rec)

    results.sort(key=lambda r: r.get("root", ""))
    print(json.dumps({"workers": workers, "candidates": len(candidates), "results": results}, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
